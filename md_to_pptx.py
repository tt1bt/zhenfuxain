
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def parse_markdown(md_content):
    slides = []
    current_slide = {"title": "", "content": []}
    
    lines = md_content.split('\n')
    for line in lines:
        line = line.strip()
        
        # 封面页检测
        if line.startswith('# '):
            if current_slide["title"]:
                slides.append(current_slide)
            current_slide = {"title": line[2:], "content": [], "type": "cover"}
            continue
        
        # 目录页检测
        if line.startswith('## 目录'):
            if current_slide["title"]:
                slides.append(current_slide)
            current_slide = {"title": "目录", "content": [], "type": "table_of_contents"}
            continue
            
        # 普通章节标题
        if line.startswith('## '):
            if current_slide["title"]:
                slides.append(current_slide)
            current_slide = {"title": line[3:], "content": [], "type": "section"}
            continue
            
        # 子标题
        if line.startswith('### '):
            current_slide["content"].append({"type": "subheading", "text": line[4:]})
            continue
            
        # 列表项
        if line.startswith('- '):
            current_slide["content"].append({"type": "bullet", "text": line[2:]})
            continue
            
        # 分隔线（新幻灯片）
        if line.startswith('---'):
            if current_slide["title"]:
                slides.append(current_slide)
            current_slide = {"title": "", "content": []}
            continue
            
        # 普通文本
        if line and not line.startswith('>') and not line.startswith('**'):
            if current_slide["content"] and current_slide["content"][-1]["type"] == "text":
                current_slide["content"][-1]["text"] += " " + line
            else:
                current_slide["content"].append({"type": "text", "text": line})
    
    if current_slide["title"] or current_slide["content"]:
        slides.append(current_slide)
    
    return slides

def create_cover_slide(prs, title):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置标题
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    
    # 设置副标题
    subtitle = slide.placeholders[1]
    subtitle.text = "毕业论文答辩\n\n姓名：[姓名]\n学号：[学号]\n指导教师：[导师姓名]\n学院/专业：[学院]/[专业]\n日期：[日期]"
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
    
    # 添加装饰形状
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(0.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 51, 102)
    shape.line.fill.background()
    
    return slide

def create_toc_slide(prs, slides):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "目录"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide.placeholders[1]
    content.text = ""
    
    idx = 1
    for s in slides:
        if s["type"] != "cover" and s["type"] != "table_of_contents":
            content.text += f"{idx}. {s['title']}\n"
            idx += 1
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    
    return slide

def create_content_slide(prs, slide_data):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = slide_data["title"]
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide.placeholders[1]
    content.text = ""
    
    for item in slide_data["content"]:
        if item["type"] == "subheading":
            content.text += f"◆ {item['text']}\n\n"
        elif item["type"] == "bullet":
            content.text += f"● {item['text']}\n"
        elif item["type"] == "text":
            content.text += f"{item['text']}\n\n"
    
    # 设置字体样式
    for i, paragraph in enumerate(content.text_frame.paragraphs):
        if paragraph.text.startswith('◆'):
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)
        elif paragraph.text.startswith('●'):
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)
        else:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)
        paragraph.alignment = PP_ALIGN.LEFT
    
    return slide

def create_thanks_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加感谢文本框
    left = Inches(2)
    top = Inches(3)
    width = Inches(6)
    height = Inches(2)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "谢谢各位老师！\n恳请批评指正！"
    
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(0, 51, 102)
    paragraph.alignment = PP_ALIGN.CENTER
    
    # 添加装饰线
    left = Inches(3)
    top = Inches(4.5)
    width = Inches(4)
    height = Inches(0.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 51, 102)
    shape.line.fill.background()
    
    return slide

def md_to_pptx(md_file_path, pptx_file_path):
    # 读取Markdown文件
    with open(md_file_path, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    # 解析Markdown
    slides = parse_markdown(md_content)
    
    # 创建PPT
    prs = Presentation()
    
    # 设置幻灯片大小为16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 添加封面页
    if slides and slides[0]["type"] == "cover":
        create_cover_slide(prs, slides[0]["title"])
        slides = slides[1:]
    
    # 添加目录页
    create_toc_slide(prs, slides)
    
    # 添加内容页
    for slide_data in slides:
        if slide_data["type"] != "table_of_contents":
            create_content_slide(prs, slide_data)
    
    # 添加致谢页
    create_thanks_slide(prs)
    
    # 保存PPT
    prs.save(pptx_file_path)
    print(f"PPT已成功生成：{pptx_file_path}")

if __name__ == "__main__":
    md_file = "毕业论文答辩大纲.md"
    pptx_file = "毕业论文答辩大纲.pptx"
    md_to_pptx(md_file, pptx_file)
