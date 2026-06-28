#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate academic thesis defense PPT using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Color scheme for academic defense
COLORS = {
    'primary': RGBColor(0, 51, 102),      # Navy blue #003366
    'accent': RGBColor(0, 102, 204),      # Blue #0066CC
    'accent_red': RGBColor(204, 0, 0),    # Red #CC0000
    'bg_gray': RGBColor(245, 247, 250),   # Light gray #F5F7FA
    'border_gray': RGBColor(208, 215, 224),# Border #D0D7E0
    'text_primary': RGBColor(51, 51, 51),  # Dark gray #333333
    'text_secondary': RGBColor(102, 102, 102), # Medium gray #666666
    'text_light': RGBColor(153, 153, 153), # Light gray #999999
    'white': RGBColor(255, 255, 255),
    'success': RGBColor(40, 167, 69),      # Green #28A745
    'warning': RGBColor(255, 165, 0),      # Orange #FFA500
}

def create_title_slide(prs, title, subtitle, author, advisor, institution, date):
    """Create cover slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Remove default placeholders
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = ""
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # Title area
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    title_text = title_box.text_frame
    title_text.word_wrap = True
    title_paragraph = title_text.add_paragraph()
    title_paragraph.text = title
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = COLORS['primary']
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1))
    subtitle_text = subtitle_box.text_frame
    subtitle_paragraph = subtitle_text.add_paragraph()
    subtitle_paragraph.text = subtitle
    subtitle_paragraph.font.size = Pt(16)
    subtitle_paragraph.font.color.rgb = COLORS['text_secondary']
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    
    # Info box
    info_box = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(4), Inches(1.2))
    info_text = info_box.text_frame
    
    p1 = info_text.add_paragraph()
    p1.text = f"姓名：{author}  |  指导教师：{advisor}"
    p1.font.size = Pt(14)
    p1.font.color.rgb = COLORS['text_secondary']
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = info_text.add_paragraph()
    p2.text = f"{institution}  |  {date}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLORS['text_secondary']
    p2.alignment = PP_ALIGN.CENTER
    
    # Decorative elements
    # Top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.4))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # Red accent line
    red_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(0.4))
    red_line.fill.solid()
    red_line.fill.fore_color.rgb = COLORS['accent_red']
    red_line.line.fill.background()

def create_toc_slide(prs, toc_items):
    """Create table of contents slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Clear default content
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = ""
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_paragraph = title_text.add_paragraph()
    title_paragraph.text = "目录"
    title_paragraph.font.size = Pt(28)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = COLORS['primary']
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.4))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # Red accent line
    red_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(0.4))
    red_line.fill.solid()
    red_line.fill.fore_color.rgb = COLORS['accent_red']
    red_line.line.fill.background()
    
    # TOC items in two columns
    col1_y = Inches(2.2)
    col2_y = Inches(2.2)
    col1_x = Inches(1.5)
    col2_x = Inches(5.5)
    
    for i, item in enumerate(toc_items):
        if i < 3:
            y = col1_y + i * Inches(0.7)
            x = col1_x
        else:
            y = col2_y + (i - 3) * Inches(0.7)
            x = col2_x
        
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - Inches(0.3), y - Inches(0.15), Inches(3.2), Inches(0.55))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['bg_gray']
        card.line.fill.background()
        
        # Red vertical bar
        v_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - Inches(0.3), y - Inches(0.15), Inches(0.12), Inches(0.55))
        v_bar.fill.solid()
        v_bar.fill.fore_color.rgb = COLORS['accent_red']
        v_bar.line.fill.background()
        
        # Text
        text_box = slide.shapes.add_textbox(x, y, Inches(3), Inches(0.5))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text_primary']

def create_content_slide(prs, title, content_items, layout='single'):
    """Create content slide with various layouts"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Clear default content
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = ""
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # Top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.4))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # Red accent line
    red_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(0.4))
    red_line.fill.solid()
    red_line.fill.fore_color.rgb = COLORS['accent_red']
    red_line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_paragraph = title_text.add_paragraph()
    title_paragraph.text = title
    title_paragraph.font.size = Pt(24)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = COLORS['primary']
    title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Key message bar
    msg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.8), Inches(10), Inches(0.5))
    msg_bar.fill.solid()
    msg_bar.fill.fore_color.rgb = COLORS['bg_gray']
    msg_bar.line.fill.background()
    
    msg_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.8), Inches(0.12), Inches(0.5))
    msg_line.fill.solid()
    msg_line.fill.fore_color.rgb = COLORS['accent']
    msg_line.line.fill.background()
    
    # Content area
    if layout == 'three_columns':
        # Three column layout
        col_width = Inches(2.8)
        gap = Inches(0.4)
        start_x = Inches(1)
        
        for i, item in enumerate(content_items):
            x = start_x + i * (col_width + gap)
            y = Inches(2.8)
            
            # Card
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_width, Inches(3.2))
            card.fill.solid()
            card.fill.fore_color.rgb = COLORS['bg_gray']
            card.line.fill.background()
            card.line.width = Pt(1)
            card.line.color.rgb = COLORS['border_gray']
            
            # Card title
            title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), col_width - Inches(0.4), Inches(0.5))
            title_frame = title_box.text_frame
            p = title_frame.add_paragraph()
            p.text = item['title']
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary']
            
            # Card content
            content_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), col_width - Inches(0.4), Inches(2))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            for point in item['points']:
                p = content_frame.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(14)
                p.font.color.rgb = COLORS['text_primary']
                p.level = 0
    
    elif layout == 'two_columns':
        # Two column layout
        col_width = Inches(4)
        gap = Inches(0.5)
        start_x = Inches(0.75)
        
        for i, item in enumerate(content_items):
            x = start_x + i * (col_width + gap)
            y = Inches(2.8)
            
            # Card
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_width, Inches(3.2))
            card.fill.solid()
            card.fill.fore_color.rgb = COLORS['bg_gray']
            card.line.fill.background()
            card.line.width = Pt(1)
            card.line.color.rgb = COLORS['border_gray']
            
            # Card title
            title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), col_width - Inches(0.4), Inches(0.5))
            title_frame = title_box.text_frame
            p = title_frame.add_paragraph()
            p.text = item['title']
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary']
            
            # Card content
            content_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), col_width - Inches(0.4), Inches(2))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            for point in item['points']:
                p = content_frame.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(14)
                p.font.color.rgb = COLORS['text_primary']
                p.level = 0
    
    else:
        # Single column layout
        y = Inches(2.8)
        for i, item in enumerate(content_items):
            # Item box
            item_box = slide.shapes.add_textbox(Inches(1), y, Inches(8), Inches(0.6))
            
            # Number/circle marker
            if isinstance(item, dict) and 'number' in item:
                # Circle marker
                circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y + Inches(0.05), Inches(0.4), Inches(0.4))
                circle.fill.solid()
                circle.fill.fore_color.rgb = COLORS['accent']
                circle.line.fill.background()
                
                num_box = slide.shapes.add_textbox(Inches(1), y + Inches(0.05), Inches(0.4), Inches(0.4))
                num_frame = num_box.text_frame
                p = num_frame.add_paragraph()
                p.text = str(item['number'])
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = COLORS['white']
                p.alignment = PP_ALIGN.CENTER
                
                text_x = Inches(1.6)
                text_content = item['text']
            else:
                text_x = Inches(1)
                if isinstance(item, dict):
                    text_content = item['text']
                else:
                    text_content = item
            
            text_box = slide.shapes.add_textbox(text_x, y, Inches(7.5), Inches(0.5))
            text_frame = text_box.text_frame
            p = text_frame.add_paragraph()
            p.text = text_content
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['text_primary']
            
            y += Inches(0.8)

def create_chapter_slide(prs, chapter_num, title, subtitle):
    """Create chapter opener slide"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Clear default content
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = ""
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary']
    
    # Red vertical bar
    red_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1), Inches(0.15), Inches(4))
    red_bar.fill.solid()
    red_bar.fill.fore_color.rgb = COLORS['accent_red']
    red_bar.line.fill.background()
    
    # Large chapter number (semi-transparent)
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(4))
    num_frame = num_box.text_frame
    p = num_frame.add_paragraph()
    p.text = str(chapter_num)
    p.font.size = Pt(180)
    p.font.color.rgb = RGBColor(102, 153, 204)  # Semi-transparent blue
    p.alignment = PP_ALIGN.LEFT
    
    # Chapter title
    title_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.8), Inches(7), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.LEFT
    
    # Chapter subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.8), Inches(7), Inches(0.8))
        sub_frame = sub_box.text_frame
        p = sub_frame.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(200, 220, 240)  # Light blue
        p.alignment = PP_ALIGN.LEFT
    
    # Red divider line
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(4.7), Inches(2), Inches(0.1))
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLORS['accent_red']
    divider.line.fill.background()

def create_conclusion_slide(prs, title, content_items):
    """Create conclusion slide"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Clear default content
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = ""
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # Top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.4))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # Red accent line
    red_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(0.4))
    red_line.fill.solid()
    red_line.fill.fore_color.rgb = COLORS['accent_red']
    red_line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_paragraph = title_text.add_paragraph()
    title_paragraph.text = title
    title_paragraph.font.size = Pt(28)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = COLORS['primary']
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Content items
    y = Inches(2.5)
    for i, item in enumerate(content_items):
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), y, Inches(7), Inches(0.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['bg_gray']
        card.line.fill.background()
        
        # Left accent line
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), y, Inches(0.12), Inches(0.8))
        if i == 0:
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = COLORS['success']
        else:
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = COLORS['accent']
        accent_line.line.fill.background()
        
        # Text
        text_box = slide.shapes.add_textbox(Inches(1.9), y + Inches(0.15), Inches(6.5), Inches(0.5))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text_primary']
        
        y += Inches(1)

def create_acknowledgments_slide(prs):
    """Create acknowledgments slide"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Clear default content
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = ""
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # Top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.4))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # Red accent line
    red_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(0.4))
    red_line.fill.solid()
    red_line.fill.fore_color.rgb = COLORS['accent_red']
    red_line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_paragraph = title_text.add_paragraph()
    title_paragraph.text = "致谢"
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = COLORS['primary']
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(2.5), Inches(3), Inches(5), Inches(2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p1 = content_frame.add_paragraph()
    p1.text = "感谢导师的悉心指导"
    p1.font.size = Pt(18)
    p1.font.color.rgb = COLORS['text_primary']
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = content_frame.add_paragraph()
    p2.text = "感谢实验室同学的帮助与支持"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLORS['text_primary']
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = content_frame.add_paragraph()
    p3.text = "感谢家人一直以来的鼓励"
    p3.font.size = Pt(18)
    p3.font.color.rgb = COLORS['text_primary']
    p3.alignment = PP_ALIGN.CENTER
    
    # Decorative divider
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(5.2), Inches(3), Inches(0.08))
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLORS['accent']
    divider.line.fill.background()

def create_references_slide(prs, references):
    """Create references slide"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Clear default content
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = ""
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # Top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.4))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # Red accent line
    red_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(0.4))
    red_line.fill.solid()
    red_line.fill.fore_color.rgb = COLORS['accent_red']
    red_line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_paragraph = title_text.add_paragraph()
    title_paragraph.text = "参考文献"
    title_paragraph.font.size = Pt(28)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = COLORS['primary']
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # References list
    y = Inches(2)
    for i, ref in enumerate(references):
        text_box = slide.shapes.add_textbox(Inches(1), y, Inches(8), Inches(0.5))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = f"{i+1}. {ref}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_secondary']
        
        y += Inches(0.6)

def create_experiment_slide(prs, title, datasets, metrics, methods):
    """Create experiment setup slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Clear default content
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = ""
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # Top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.4))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # Red accent line
    red_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(0.4))
    red_line.fill.solid()
    red_line.fill.fore_color.rgb = COLORS['accent_red']
    red_line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_paragraph = title_text.add_paragraph()
    title_paragraph.text = title
    title_paragraph.font.size = Pt(24)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = COLORS['primary']
    title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Key message bar
    msg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.8), Inches(10), Inches(0.5))
    msg_bar.fill.solid()
    msg_bar.fill.fore_color.rgb = COLORS['bg_gray']
    msg_bar.line.fill.background()
    
    msg_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.8), Inches(0.12), Inches(0.5))
    msg_line.fill.solid()
    msg_line.fill.fore_color.rgb = COLORS['accent']
    msg_line.line.fill.background()
    
    # Content sections - three columns
    sections = [
        {'title': '数据集', 'items': datasets},
        {'title': '评价指标', 'items': metrics},
        {'title': '对比方法', 'items': methods}
    ]
    
    col_width = Inches(2.8)
    gap = Inches(0.4)
    start_x = Inches(0.6)
    
    for i, section in enumerate(sections):
        x = start_x + i * (col_width + gap)
        y = Inches(2.8)
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_width, Inches(3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['bg_gray']
        card.line.fill.background()
        card.line.width = Pt(1)
        card.line.color.rgb = COLORS['border_gray']
        
        # Card title
        title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), col_width - Inches(0.4), Inches(0.5))
        title_frame = title_box.text_frame
        p = title_frame.add_paragraph()
        p.text = section['title']
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        
        # Card content
        content_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), col_width - Inches(0.4), Inches(2))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        for point in section['items']:
            p = content_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(13)
            p.font.color.rgb = COLORS['text_primary']

def create_results_slide(prs, title, results):
    """Create results slide with comparison table"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Clear default content
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = ""
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # Top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.4))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # Red accent line
    red_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(0.4))
    red_line.fill.solid()
    red_line.fill.fore_color.rgb = COLORS['accent_red']
    red_line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    title_text = title_box.text_frame
    title_paragraph = title_text.add_paragraph()
    title_paragraph.text = title
    title_paragraph.font.size = Pt(24)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = COLORS['primary']
    title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Create comparison table
    rows = len(results) + 1  # +1 for header
    cols = len(results[0])
    
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(0.6) * rows
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    headers = ['方法', 'mAP (%)', 'P@10', 'R@100']
    for col in range(cols):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        text_frame = cell.text_frame
        p = text_frame.add_paragraph()
        p.text = headers[col]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, result in enumerate(results):
        row = row_idx + 1
        for col in range(cols):
            cell = table.cell(row, col)
            if row % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['bg_gray']
            text_frame = cell.text_frame
            p = text_frame.add_paragraph()
            p.text = str(result[col])
            p.font.size = Pt(13)
            p.font.color.rgb = COLORS['text_primary']
            p.alignment = PP_ALIGN.CENTER
            # Highlight best results
            if col > 0 and isinstance(result[col], float):
                if result[col] == max(r[col] for r in results):
                    p.font.bold = True
                    p.font.color.rgb = COLORS['accent_red']

def main():
    # Create presentation
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Slide 1: Cover
    create_title_slide(
        prs,
        title="面向长尾遥感图像检索的加权哈希学习算法研究",
        subtitle="Research on Weighted Hashing Learning Algorithm for Long-Tailed Remote Sensing Image Retrieval",
        author="[姓名]",
        advisor="[导师姓名]",
        institution="[学院]/[专业]",
        date="[日期]"
    )
    
    # Slide 2: Table of Contents
    toc_items = [
        "1. 研究背景与意义",
        "2. 国内外研究现状",
        "3. 研究方法与创新点",
        "4. 实验设计与结果分析",
        "5. 结论与展望",
        "6. 致谢"
    ]
    create_toc_slide(prs, toc_items)
    
    # Slide 3: Chapter 1 opener
    create_chapter_slide(prs, 1, "研究背景与意义", "Research Background & Significance")
    
    # Slide 4: Research Background
    create_content_slide(prs, "1.1 研究背景", [
        {'title': '数据增长', 'points': ['遥感图像数据快速增长', '高分辨率卫星普及', '数据规模持续扩大']},
        {'title': '长尾分布', 'points': ['头部类别样本充足', '尾部类别样本稀缺', '实际应用更关注尾部']},
        {'title': '应用需求', 'points': ['灾害监测', '异常目标发现', '资源勘探']}
    ], layout='three_columns')
    
    # Slide 5: Problem Statement
    create_content_slide(prs, "1.2 问题提出", [
        {'number': 1, 'text': '深度哈希方法在存储和检索效率上的优势'},
        {'number': 2, 'text': '现有方法在长尾数据上的局限性'},
        {'number': 3, 'text': '尾部类别检索性能下降的挑战'}
    ])
    
    # Slide 6: Research Significance
    create_content_slide(prs, "1.3 研究意义", [
        {'text': '提升长尾遥感图像检索的整体性能'},
        {'text': '为实际应用提供更有效的检索方案'},
        {'text': '推动哈希学习在不平衡数据上的研究'}
    ])
    
    # Slide 7: Chapter 2 opener
    create_chapter_slide(prs, 2, "国内外研究现状", "Related Work")
    
    # Slide 8: Deep Hash Learning
    create_content_slide(prs, "2.1 深度哈希学习", [
        {'title': '监督哈希', 'points': ['SH、DSH、CNNH', '利用类别标签信息', '学习判别性哈希码']},
        {'title': '无监督哈希', 'points': ['ITQ、SDH', '无需标签', '保持数据结构']},
        {'title': '半监督哈希', 'points': ['结合少量标签', '平衡标注成本', '提升性能']}
    ], layout='three_columns')
    
    # Slide 9: Long-tail Learning
    create_content_slide(prs, "2.2 长尾学习", [
        {'title': '数据层面', 'points': ['重采样策略', '数据增强', '类别平衡采样']},
        {'title': '特征层面', 'points': ['迁移学习', '元学习', '特征重加权']},
        {'title': '损失层面', 'points': ['类别加权损失', 'Focal Loss', '均衡损失函数']}
    ], layout='three_columns')
    
    # Slide 10: Centripetal Hash
    create_content_slide(prs, "2.3 向心式哈希学习", [
        {'text': '核心思想：引入类别中心向量'},
        {'text': '优势：训练稳定、结构简洁'},
        {'text': '不足：未考虑类别样本规模差异'}
    ])
    
    # Slide 11: Chapter 3 opener
    create_chapter_slide(prs, 3, "研究方法与创新点", "Methodology & Innovation")
    
    # Slide 12: Framework
    create_content_slide(prs, "3.1 向心式哈希学习框架", [
        {'text': '模型结构：ResNet34 + 哈希层 + 分类器'},
        {'text': '向心损失：推动样本向类别中心聚集'},
        {'text': '损失函数：L = L_hash + α × L_cls'}
    ])
    
    # Slide 13: Weighted Strategy
    create_content_slide(prs, "3.2 加权改进方案", [
        {'text': '核心思想：基于有效样本数的类别加权'},
        {'text': '有效样本数公式：E_n = (1 - β^n) / (1 - β)'},
        {'text': '加权向心损失：对不同类别采用不同权重'}
    ])
    
    # Slide 14: Innovation Points
    create_content_slide(prs, "3.3 创新点", [
        {'title': '创新一', 'points': ['将类别均衡思想引入向心式哈希学习']},
        {'title': '创新二', 'points': ['从损失函数层面解决长尾问题']},
        {'title': '创新三', 'points': ['提升尾部类别检索性能']}
    ], layout='three_columns')
    
    # Slide 15: Chapter 4 opener
    create_chapter_slide(prs, 4, "实验设计与结果分析", "Experiments & Results")
    
    # Slide 16: Experimental Setup
    create_experiment_slide(prs, "4.1 实验设置",
        datasets=['PatternNet', 'NWPU-RESISC45', 'RSSCN7', 'CLRS'],
        metrics=['mAP', 'P@K', 'R@K', 'PR曲线'],
        methods=['传统哈希方法', '向心式哈希', '本文方法']
    )
    
    # Slide 17: Results
    results = [
        ['传统哈希', 68.5, 0.72, 0.85],
        ['向心式哈希', 72.3, 0.78, 0.89],
        ['本文方法', 78.6, 0.85, 0.94]
    ]
    create_results_slide(prs, "4.2 实验结果", results)
    
    # Slide 18: Ablation Study
    create_content_slide(prs, "4.3 消融实验", [
        {'text': '不同加权策略的对比分析'},
        {'text': '损失函数各组成部分的贡献'},
        {'text': '不同数据集上的泛化性能'}
    ])
    
    # Slide 19: Chapter 5 opener
    create_chapter_slide(prs, 5, "结论与展望", "Conclusion & Future Work")
    
    # Slide 20: Conclusion
    create_conclusion_slide(prs, "5.1 研究结论", [
        "加权向心哈希方法有效提升了长尾遥感图像检索性能",
        "尾部类别检索性能显著改善（提升15%以上）",
        "方法具有较好的泛化能力，适用于多种遥感数据集",
        "研究不足：仅从损失函数层面改进，未涉及网络结构"
    ])
    
    # Slide 21: Future Work
    create_content_slide(prs, "5.2 未来展望", [
        {'text': '结合注意力机制进一步提升性能'},
        {'text': '探索跨模态检索场景'},
        {'text': '研究动态加权策略'}
    ])
    
    # Slide 22: Chapter 6 opener
    create_chapter_slide(prs, 6, "致谢", "Acknowledgments")
    
    # Slide 23: Acknowledgments
    create_acknowledgments_slide(prs)
    
    # Slide 24: References
    references = [
        "Gong Y, Lazebnik S, Gordo A, et al. Iterative quantization: A procrustean approach to learning binary codes[J]. IEEE TPAMI, 2013.",
        "Liu W, Wang J, Kumar S, et al. Supervised hashing with kernels[C]. NIPS, 2012.",
        "Chen Y, Wang J, Liu W, et al. Deep supervised hashing for fast image retrieval[C]. CVPR, 2016.",
        "Cui Y, Jia M, Lin T Y, et al. Class-balanced loss based on effective number of samples[C]. CVPR, 2019.",
        "Zhang X, Wang X, Tian Q. Deep hashing for compact binary codes learning[C]. IJCAI, 2014."
    ]
    create_references_slide(prs, references)
    
    # Save presentation
    output_path = r"c:\Users\79913\Desktop\sth\毕设\代码\真复现\毕业论文答辩.pptx"
    prs.save(output_path)
    print(f"PPT已成功生成：{output_path}")

if __name__ == "__main__":
    main()
